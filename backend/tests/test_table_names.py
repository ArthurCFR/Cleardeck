"""Regression tests for stacked names in tables (Prénom / NOM matrices).

Run standalone (no pytest needed):  python backend/tests/test_table_names.py

These tests do NOT load CamemBERT (~400 MB). Instead they fake the model's
*gluing* behaviour — with aggregation_strategy="simple" CamemBERT merges
contiguous name tokens across newlines into one giant PER span — and then
exercise the real extraction -> split -> merge -> replace -> deanonymize chain.
"""

from __future__ import annotations

import io
import sys
from pathlib import Path

# Make `backend` importable when run directly from the repo root or elsewhere.
sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from pptx import Presentation
from pptx.util import Inches

from backend.engine import ai_detector, anonymizer, deanonymizer
from backend.engine.ai_detector import _split_on_newlines
from backend.engine.entity_merger import merge_entities


# --------------------------------------------------------------------------- #
# Fake NER pipeline — reproduces the cross-line gluing bug.
# --------------------------------------------------------------------------- #
def _looks_like_name(line: str) -> bool:
    tokens = line.split()
    if not tokens or len(tokens) > 3:
        return False
    return all(t[0].isupper() for t in tokens if t and t[0].isalpha())


def fake_ner(text: str):
    """Emit one PER span per *maximal run* of consecutive name-like lines,
    spanning across the newlines — exactly what real CamemBERT does."""
    out = []
    offset = 0
    run = None  # [start_char, end_char]
    spans = []
    for line in text.split("\n"):
        stripped = line.strip()
        if stripped and _looks_like_name(stripped):
            lead = len(line) - len(line.lstrip())
            s = offset + lead
            e = offset + len(line.rstrip())
            run = [s, e] if run is None else [run[0], e]
        elif run is not None:
            spans.append(tuple(run))
            run = None
        offset += len(line) + 1  # +1 for the consumed "\n"
    if run is not None:
        spans.append(tuple(run))

    for s, e in spans:
        out.append({"entity_group": "PER", "word": text[s:e],
                    "start": s, "end": e, "score": 0.999})
    return out


def _patch_pipeline():
    ai_detector.get_pipeline = lambda: fake_ner  # type: ignore[assignment]


# --------------------------------------------------------------------------- #
# Tests
# --------------------------------------------------------------------------- #
def test_split_on_newlines():
    """A glued multi-line span is broken into accurate per-line fragments."""
    text = "Hélène\nETTAHAR\nPierre CAYET"
    glued = [{"entity": text, "label": "PER", "start": 0,
              "end": len(text), "confidence": 0.999}]
    frags = _split_on_newlines(glued)
    assert [f["entity"] for f in frags] == ["Hélène", "ETTAHAR", "Pierre CAYET"]
    # Offsets must point back at the real substring.
    for f in frags:
        assert text[f["start"]:f["end"]] == f["entity"]
    # A single-line span passes through untouched.
    single = [{"entity": "Jean Dupont", "label": "PER", "start": 5,
               "end": 16, "confidence": 0.9}]
    assert _split_on_newlines(single) == single
    print("OK  test_split_on_newlines")


def test_detect_entities_no_glue():
    """Stacked names must come out as separate entities, never one mega-span."""
    _patch_pipeline()
    text = "Hélène\nETTAHAR\nPierre CAYET\nEmmanuelle MAILLOT"
    ents = [e["entity"] for e in ai_detector.detect_entities(text)]
    assert ents == ["Hélène", "ETTAHAR", "Pierre CAYET", "Emmanuelle MAILLOT"], ents
    print("OK  test_detect_entities_no_glue")


def test_merge_no_collapse():
    """Distinct fragments get distinct placeholders; identical ones share."""
    ai = [{"entity": n, "category": "personnes"} for n in
          ["Hélène", "ETTAHAR", "Pierre CAYET", "Hélène"]]  # last is a dup
    pairs, entries = merge_entities(project_entities=None, ai_confirmed=ai)
    ph = {orig.lower(): p for orig, p in pairs}
    assert ph["hélène"] != ph["ettahar"] != ph["pierre cayet"]
    assert ph["hélène"] != ph["pierre cayet"]
    # Each placeholder maps to exactly one distinct original (de-anon safe).
    by_ph: dict[str, set[str]] = {}
    for e in entries:
        by_ph.setdefault(e["placeholder"], set()).add(e["original"].lower())
    assert all(len(v) == 1 for v in by_ph.values()), by_ph
    print("OK  test_merge_no_collapse")


def _build_pptx() -> bytes:
    """Table whose cells stack names across paragraphs."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    rows, cols = 2, 1
    tbl = slide.shapes.add_table(
        rows, cols, Inches(1), Inches(1), Inches(4), Inches(2)).table

    def fill(cell, lines):
        tf = cell.text_frame
        tf.paragraphs[0].text = lines[0]
        for ln in lines[1:]:
            tf.add_paragraph().text = ln

    fill(tbl.cell(0, 0), ["Hélène", "ETTAHAR"])              # 1 person, 2 lines
    fill(tbl.cell(1, 0), ["Pierre CAYET", "Emmanuelle MAILLOT"])  # 2 people

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_pptx_round_trip():
    _patch_pipeline()
    original = _build_pptx()

    out_bytes, mapping, _ = anonymizer.anonymize(
        original, "pptx", "matrice.pptx", project_id=None)

    # Every distinct name must own its own placeholder.
    placeholders = {e["placeholder"] for e in mapping["entities"]}
    assert len(placeholders) == 4, mapping["entities"]
    # No placeholder is ambiguous (one placeholder -> one original).
    by_ph: dict[str, set[str]] = {}
    for e in mapping["entities"]:
        by_ph.setdefault(e["placeholder"], set()).add(e["original"])
    assert all(len(v) == 1 for v in by_ph.values()), by_ph

    # Anonymised output must no longer contain any of the names.
    anon_text = _all_text(out_bytes)
    for name in ["Hélène", "ETTAHAR", "Pierre CAYET", "Emmanuelle MAILLOT"]:
        assert name not in anon_text, f"leak: {name}"

    # Round-trip: deanonymize restores the original text exactly.
    restored = deanonymizer.deanonymize(out_bytes, "pptx", mapping)
    assert _all_text(restored) == _all_text(original)
    print("OK  test_pptx_round_trip")


def _all_text(pptx_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(pptx_bytes))
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for p in cell.text_frame.paragraphs:
                            out.append(p.text)
    return "\n".join(out)


if __name__ == "__main__":
    test_split_on_newlines()
    test_detect_entities_no_glue()
    test_merge_no_collapse()
    test_pptx_round_trip()
    print("\nAll tests passed.")
