"""Regression test: PPTX paragraph XML order after replacement.

Run standalone (no pytest needed):
    python backend/tests/test_pptx_paragraph_order.py

Bug: when a paragraph carries an <a:endParaRPr> (common in real decks) and one
of its runs is replaced, the rebuilt <a:r> elements were appended *after*
<a:endParaRPr>. The OOXML schema requires <a:endParaRPr> to be the LAST child of
<a:p>. python-pptx reads the text anyway (lenient), but PowerPoint is strict and
silently drops the whole paragraph/shape — the block "disappears" with no
placeholder, more often when many entities hit the same block. The fix inserts
new runs before <a:endParaRPr>.
"""

from __future__ import annotations

import io
import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parents[2]))

from lxml import etree
from pptx import Presentation
from pptx.util import Inches

from backend.engine import pptx_handler

A = "http://schemas.openxmlformats.org/drawingml/2006/main"


def _q(tag: str) -> str:
    return f"{{{A}}}{tag}"


def _build_deck_paragraph(runs_text, with_pPr=True, with_endParaRPr=True) -> Presentation:
    """A textbox paragraph mirroring a real slide: optional pPr + endParaRPr."""
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tf = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2)).text_frame
    p = tf.paragraphs[0]
    for t in runs_text:
        r = p.add_run()
        r.text = t
    para = p._p
    if with_pPr:
        para.insert(0, etree.SubElement(para, _q("pPr")))
    if with_endParaRPr:
        etree.SubElement(para, _q("endParaRPr"))
    return prs


def _first_para_order(prs: Presentation):
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                return [etree.QName(c).localname for c in sh.text_frame.paragraphs[0]._p]
    return []


def _reopen_text(prs: Presentation) -> str:
    buf = io.BytesIO()
    prs.save(buf)
    pr2 = Presentation(io.BytesIO(buf.getvalue()))
    out = []
    for s in pr2.slides:
        for sh in s.shapes:
            if sh.has_text_frame:
                for pp in sh.text_frame.paragraphs:
                    out.append(pp.text)
    return "\n".join(out)


def test_endpararpr_stays_last():
    """endParaRPr must remain the final child; pPr must stay first."""
    prs = _build_deck_paragraph(
        ["La ", "Direction Juridique", " gère ", "Paris", " et ", "Lille"]
    )
    pptx_handler.apply_replacements(
        prs,
        [("Direction Juridique", "[ENTREPRISE_1]"), ("Paris", "[LIEU_2]"), ("Lille", "[LIEU_1]")],
        alt_mapping={},
    )
    order = _first_para_order(prs)
    assert order[0] == "pPr", f"pPr must be first, got {order}"
    assert order[-1] == "endParaRPr", f"endParaRPr must be last, got {order}"
    # No run leaks after endParaRPr.
    assert order.index("endParaRPr") == len(order) - 1, order
    assert _reopen_text(prs) == "La [ENTREPRISE_1] gère [LIEU_2] et [LIEU_1]"
    print("OK  test_endpararpr_stays_last")


def test_no_endpararpr_still_works():
    """Paragraphs without endParaRPr keep working (append path)."""
    prs = _build_deck_paragraph(["La ", "Direction Juridique", " ok"], with_endParaRPr=False)
    pptx_handler.apply_replacements(prs, [("Direction Juridique", "[ENTREPRISE_1]")], alt_mapping={})
    order = _first_para_order(prs)
    assert "endParaRPr" not in order, order
    assert _reopen_text(prs) == "La [ENTREPRISE_1] ok"
    print("OK  test_no_endpararpr_still_works")


if __name__ == "__main__":
    test_endpararpr_stays_last()
    test_no_endpararpr_still_works()
    print("\nAll tests passed.")
