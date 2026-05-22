"""
Image handler — extract, compare via phash, and replace images in DOCX/PPTX.

Three-layer extraction for PPTX: Slide Masters, Slide Layouts, Individual Slides.
For DOCX: inline images from paragraphs and headers/footers.

Logo matching uses imagehash (perceptual hash) with two thresholds:
  - hash_diff <= 8:  AUTO_REMOVE (automatic, silent)
  - hash_diff <= 18: REVIEW (proposed to user)
  - hash_diff > 18:  KEEP (untouched)
"""

from __future__ import annotations

import io
import base64
from dataclasses import dataclass

import imagehash
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from docx import Document
from docx.shared import Inches
from lxml import etree


THRESHOLD_AUTO = 8
THRESHOLD_REVIEW = 18

REPLACEMENT_TEXT = "[IMAGE SUPPRIMEE]"
REPLACEMENT_COLOR = (200, 200, 200)
REPLACEMENT_TEXT_COLOR = (100, 100, 100)


@dataclass
class ExtractedImage:
    """An image extracted from a document with its metadata."""
    image_bytes: bytes
    location: str           # e.g. "Slide Master 1", "Slide 3", "En-tete"
    width: int              # pixels (estimated from EMU or inline extent)
    height: int
    phash: str              # hex string of perceptual hash
    thumbnail_b64: str      # base64-encoded JPEG thumbnail for UI
    shape_ref: object       # reference to shape/inline element for replacement
    source_layer: str       # "master", "layout", "slide", "docx_inline"
    # Set after comparison:
    status: str = "keep"    # "auto_remove", "review", "keep"
    hash_diff: int | None = None


def compute_phash(image_bytes: bytes) -> str:
    """Compute perceptual hash of an image.

    RGBA images are composited onto a white background before hashing
    so that transparent PNGs produce the same hash as their opaque
    rendering inside PPTX/DOCX documents.
    """
    img = Image.open(io.BytesIO(image_bytes))
    if img.mode == "RGBA":
        background = Image.new("RGB", img.size, (255, 255, 255))
        background.paste(img, mask=img.split()[3])
        img = background
    return str(imagehash.phash(img))


def _make_thumbnail_b64(image_bytes: bytes, max_size: int = 150) -> str:
    """Create a base64-encoded JPEG thumbnail."""
    img = Image.open(io.BytesIO(image_bytes))
    img.thumbnail((max_size, max_size))
    if img.mode in ("RGBA", "P"):
        img = img.convert("RGB")
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=75)
    return base64.b64encode(buf.getvalue()).decode("ascii")


def _generate_grey_replacement(width: int, height: int) -> bytes:
    """Generate a grey rectangle PNG with centered "[IMAGE SUPPRIMEE]" text."""
    # Ensure minimum dimensions
    w = max(width, 100)
    h = max(height, 60)
    img = Image.new("RGB", (w, h), REPLACEMENT_COLOR)
    draw = ImageDraw.Draw(img)

    # Try to use a reasonable font size relative to image
    font_size = max(12, min(w // 20, h // 4, 36))
    try:
        font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", font_size)
    except (OSError, IOError):
        try:
            font = ImageFont.truetype("/usr/share/fonts/TTF/DejaVuSans.ttf", font_size)
        except (OSError, IOError):
            font = ImageFont.load_default()

    bbox = draw.textbbox((0, 0), REPLACEMENT_TEXT, font=font)
    text_w = bbox[2] - bbox[0]
    text_h = bbox[3] - bbox[1]
    x = (w - text_w) // 2
    y = (h - text_h) // 2
    draw.text((x, y), REPLACEMENT_TEXT, fill=REPLACEMENT_TEXT_COLOR, font=font)

    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


# ============================================================
# PPTX image extraction (3 layers)
# ============================================================

def _extract_images_from_shapes(shapes, location_prefix: str, source_layer: str) -> list[ExtractedImage]:
    """Extract images from a collection of shapes."""
    results = []
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                image = shape.image
                image_bytes = image.blob
                img = Image.open(io.BytesIO(image_bytes))
                w, h = img.size

                results.append(ExtractedImage(
                    image_bytes=image_bytes,
                    location=location_prefix,
                    width=w,
                    height=h,
                    phash=compute_phash(image_bytes),
                    thumbnail_b64=_make_thumbnail_b64(image_bytes),
                    shape_ref=shape,
                    source_layer=source_layer,
                ))
            except Exception:
                continue

        # Check group shapes recursively
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                results.extend(_extract_images_from_shapes(
                    shape.shapes, location_prefix, source_layer
                ))
            except Exception:
                continue
    return results


def extract_images_pptx(prs: Presentation) -> list[ExtractedImage]:
    """Extract all images from a PPTX across 3 layers."""
    images = []

    # Layer 1: Slide Masters
    for i, master in enumerate(prs.slide_masters):
        label = f"Slide Master {i + 1}"
        images.extend(_extract_images_from_shapes(master.shapes, label, "master"))

    # Layer 2: Slide Layouts
    for master in prs.slide_masters:
        for j, layout in enumerate(master.slide_layouts):
            label = f"Layout — {layout.name or f'Layout {j + 1}'}"
            images.extend(_extract_images_from_shapes(layout.shapes, label, "layout"))

    # Layer 3: Individual Slides
    for k, slide in enumerate(prs.slides):
        label = f"Slide {k + 1}"
        images.extend(_extract_images_from_shapes(slide.shapes, label, "slide"))

    return images


# ============================================================
# DOCX image extraction
# ============================================================

def extract_images_docx(doc: Document) -> list[ExtractedImage]:
    """Extract inline images from a DOCX document."""
    images = []
    nsmap = {
        'wp': 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'pic': 'http://schemas.openxmlformats.org/drawingml/2006/picture',
    }

    def _extract_from_paragraphs(paragraphs, location_prefix: str):
        for para in paragraphs:
            for run in para.runs:
                # Find inline images (drawings)
                drawings = run._r.findall('.//wp:inline', nsmap)
                drawings += run._r.findall('.//wp:anchor', nsmap)
                for drawing in drawings:
                    try:
                        blip = drawing.find('.//a:blip', nsmap)
                        if blip is None:
                            continue
                        r_embed = blip.get(f'{{{nsmap["r"]}}}embed')
                        if not r_embed:
                            continue

                        # Get the image part
                        rel = para.part.rels.get(r_embed)
                        if rel is None:
                            continue
                        image_part = rel.target_part
                        image_bytes = image_part.blob

                        img = Image.open(io.BytesIO(image_bytes))
                        w, h = img.size

                        # Get extent for replacement sizing (in EMU)
                        extent = drawing.find('.//wp:extent', nsmap)
                        if extent is None:
                            extent = drawing.find('wp:extent', nsmap)

                        images.append(ExtractedImage(
                            image_bytes=image_bytes,
                            location=location_prefix,
                            width=w,
                            height=h,
                            phash=compute_phash(image_bytes),
                            thumbnail_b64=_make_thumbnail_b64(image_bytes),
                            shape_ref={
                                'run': run,
                                'drawing': drawing,
                                'blip': blip,
                                'r_embed': r_embed,
                                'para': para,
                            },
                            source_layer="docx_inline",
                        ))
                    except Exception:
                        continue

    # Body paragraphs
    _extract_from_paragraphs(doc.paragraphs, "Corps du document")

    # Tables
    for t_idx, table in enumerate(doc.tables):
        for row in table.rows:
            for cell in row.cells:
                _extract_from_paragraphs(cell.paragraphs, f"Tableau {t_idx + 1}")

    # Headers and footers
    for section in doc.sections:
        if section.header:
            _extract_from_paragraphs(section.header.paragraphs, "En-tete")
        if section.footer:
            _extract_from_paragraphs(section.footer.paragraphs, "Pied de page")

    return images


# ============================================================
# Logo matching
# ============================================================

def match_against_logos(
    images: list[ExtractedImage],
    logo_hashes: list[str],
) -> list[ExtractedImage]:
    """Compare extracted images against stored logo hashes.

    Updates each image's status and hash_diff fields in-place.
    """
    if not logo_hashes:
        # No project logos — all images go to review
        for img in images:
            img.status = "review"
            img.hash_diff = None
        return images

    for img in images:
        img_hash = imagehash.hex_to_hash(img.phash)
        min_diff = min(
            img_hash - imagehash.hex_to_hash(lh)
            for lh in logo_hashes
        )
        img.hash_diff = min_diff

        if min_diff <= THRESHOLD_AUTO:
            img.status = "auto_remove"
        elif min_diff <= THRESHOLD_REVIEW:
            img.status = "review"
        else:
            img.status = "keep"

    return images


# ============================================================
# Image replacement
# ============================================================

def apply_image_removals_pptx(prs: Presentation, images_to_remove: list[ExtractedImage]):
    """Replace marked images with grey rectangles in a PPTX.

    For each image to remove, we replace the picture's image blob with
    a generated grey PNG of the same dimensions.
    """
    for img in images_to_remove:
        shape = img.shape_ref
        if shape is None:
            continue
        try:
            replacement_png = _generate_grey_replacement(img.width, img.height)

            # Access the image part via the blip relationship
            slide_part = shape.part
            blip = shape._element.find('.//' + _pptx_ns('a:blip'))
            if blip is None:
                continue
            r_id = blip.get(
                '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed'
            )
            if r_id:
                rel = slide_part.rels[r_id]
                rel.target_part._blob = replacement_png
        except Exception:
            # Fallback: remove the shape entirely
            try:
                sp = shape._element
                sp.getparent().remove(sp)
            except Exception:
                pass


def apply_image_removals_docx(doc: Document, images_to_remove: list[ExtractedImage]):
    """Replace marked images with grey rectangles in a DOCX."""
    for img in images_to_remove:
        ref = img.shape_ref
        if ref is None or not isinstance(ref, dict):
            continue
        try:
            replacement_png = _generate_grey_replacement(img.width, img.height)

            r_embed = ref['r_embed']
            para = ref['para']

            rel = para.part.rels.get(r_embed)
            if rel:
                rel.target_part._blob = replacement_png
        except Exception:
            # If replacement fails, try to remove the drawing element
            try:
                drawing = ref['drawing']
                drawing.getparent().remove(drawing)
            except Exception:
                pass


def _pptx_ns(tag: str) -> str:
    """Resolve a prefixed tag to its full namespace form for pptx XML."""
    ns_map = {
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
        'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    }
    prefix, local = tag.split(':')
    return f'{{{ns_map[prefix]}}}{local}'


# ============================================================
# Public API for the pipeline
# ============================================================

def preview_images(
    file_bytes: bytes,
    file_type: str,
    logo_hashes: list[str] | None = None,
) -> dict:
    """Extract images from a file and match against project logos.

    Returns a preview dict for the frontend.
    """
    if file_type == "pptx":
        prs = Presentation(io.BytesIO(file_bytes))
        images = extract_images_pptx(prs)
    elif file_type == "docx":
        doc = Document(io.BytesIO(file_bytes))
        images = extract_images_docx(doc)
    else:
        return {"images": [], "auto_removed_count": 0, "review_count": 0}

    match_against_logos(images, logo_hashes or [])

    auto_removed = [img for img in images if img.status == "auto_remove"]
    review = [img for img in images if img.status == "review"]

    return {
        "images": [
            {
                "index": i,
                "location": img.location,
                "width": int(img.width),
                "height": int(img.height),
                "thumbnail_b64": img.thumbnail_b64,
                "status": img.status,
                "hash_diff": int(img.hash_diff) if img.hash_diff is not None else None,
                "source_layer": img.source_layer,
            }
            for i, img in enumerate(images)
        ],
        "auto_removed_count": len(auto_removed),
        "review_count": len(review),
    }


def apply_image_anonymization(
    file_bytes: bytes,
    file_type: str,
    logo_hashes: list[str] | None = None,
    user_remove_indices: list[int] | None = None,
) -> bytes:
    """Apply image removals to a file.

    Args:
        file_bytes: raw file bytes
        file_type: "docx" or "pptx"
        logo_hashes: project logo hashes for auto-matching
        user_remove_indices: indices of images the user confirmed for removal

    Returns:
        modified file bytes
    """
    user_remove_indices = user_remove_indices or []

    if file_type == "pptx":
        prs = Presentation(io.BytesIO(file_bytes))
        images = extract_images_pptx(prs)
        match_against_logos(images, logo_hashes or [])

        to_remove = [
            img for i, img in enumerate(images)
            if img.status == "auto_remove" or i in user_remove_indices
        ]
        apply_image_removals_pptx(prs, to_remove)

        buf = io.BytesIO()
        prs.save(buf)
        return buf.getvalue()

    elif file_type == "docx":
        doc = Document(io.BytesIO(file_bytes))
        images = extract_images_docx(doc)
        match_against_logos(images, logo_hashes or [])

        to_remove = [
            img for i, img in enumerate(images)
            if img.status == "auto_remove" or i in user_remove_indices
        ]
        apply_image_removals_docx(doc, to_remove)

        buf = io.BytesIO()
        doc.save(buf)
        return buf.getvalue()

    return file_bytes
