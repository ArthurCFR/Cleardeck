"""
PPTX file handler — extract text with run-level metadata and apply replacements.

Handles: slide text frames, table cells, slide notes, and grouped shapes (recursive).
"""

from __future__ import annotations

import io
from dataclasses import dataclass, field
from lxml import etree
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from .cross_run import replace_entities_in_paragraph_pptx, replace_entities_in_text


# TODO: Images containing text (would need OCR)
# TODO: Charts and SmartArt in PPTX


@dataclass
class TextBlock:
    """A block of text extracted from a PPTX file with its location."""
    section: str  # e.g. "Slide 1 — Titre", "Slide 3 — Notes"
    text: str
    paragraph_ref: object = field(default=None, repr=False)


def _extract_from_shapes(shapes, slide_label: str, blocks: list[TextBlock]):
    """Recursively extract text from shapes, including grouped shapes."""
    for shape in shapes:
        # Recurse into group shapes
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                _extract_from_shapes(shape.shapes, slide_label, blocks)
            except Exception:
                pass
            continue

        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    shape_name = shape.name or "Forme"
                    blocks.append(TextBlock(
                        section=f"{slide_label} — {shape_name}",
                        text=text,
                        paragraph_ref=para,
                    ))

        if shape.has_table:
            table = shape.table
            for r_idx, row in enumerate(table.rows):
                for c_idx, cell in enumerate(row.cells):
                    for para in cell.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            blocks.append(TextBlock(
                                section=f"{slide_label} — Tableau, Cellule ({r_idx + 1},{c_idx + 1})",
                                text=text,
                                paragraph_ref=para,
                            ))


def extract_text_blocks(file_bytes: bytes) -> tuple[Presentation, list[TextBlock]]:
    """Extract all text blocks from a PPTX file.

    Returns:
        prs: the Presentation object (for later modification)
        blocks: list of TextBlock with text and location metadata
    """
    prs = Presentation(io.BytesIO(file_bytes))
    blocks = []

    for slide_idx, slide in enumerate(prs.slides):
        slide_label = f"Slide {slide_idx + 1}"
        _extract_from_shapes(slide.shapes, slide_label, blocks)

        # Slide notes
        if slide.has_notes_slide and slide.notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf:
                for para in notes_tf.paragraphs:
                    text = para.text.strip()
                    if text:
                        blocks.append(TextBlock(
                            section=f"{slide_label} — Notes",
                            text=text,
                            paragraph_ref=para,
                        ))

    return prs, blocks


def _replace_in_shapes(shapes, entities_sorted: list[tuple[str, str]]):
    """Recursively apply replacements in shapes, including grouped shapes."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                _replace_in_shapes(shape.shapes, entities_sorted)
            except Exception:
                pass
            continue

        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                replace_entities_in_paragraph_pptx(para, entities_sorted)

        if shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    for para in cell.text_frame.paragraphs:
                        replace_entities_in_paragraph_pptx(para, entities_sorted)


def apply_replacements(
    prs: Presentation,
    entities_sorted: list[tuple[str, str]],
    alt_mapping: dict[str, str] | None = None,
):
    """Apply entity replacements across the entire PPTX presentation.

    Args:
        prs: python-pptx Presentation object
        entities_sorted: list of (original, placeholder) sorted by len DESC
        alt_mapping: dict of original_alt_text -> placeholder for alt text replacement
    """
    for slide in prs.slides:
        _replace_in_shapes(slide.shapes, entities_sorted)

        # Slide notes
        if slide.has_notes_slide and slide.notes_slide:
            notes_tf = slide.notes_slide.notes_text_frame
            if notes_tf:
                for para in notes_tf.paragraphs:
                    replace_entities_in_paragraph_pptx(para, entities_sorted)

    # Replace alt texts
    if alt_mapping:
        _replace_alt_texts_in_shapes_all(prs, alt_mapping)

    # Anonymize metadata (title gets entity replacement, rest is cleared)
    _clean_metadata(prs, entities_sorted)


def extract_alt_texts(prs: Presentation) -> list[tuple[str, str]]:
    """Extract all non-empty alt texts from shapes across the presentation.

    Returns:
        list of (alt_text, location_label)
    """
    alt_texts: list[tuple[str, str]] = []

    for slide_idx, slide in enumerate(prs.slides):
        _collect_alt_texts(slide.shapes, f"Slide {slide_idx + 1}", alt_texts)

    for i, master in enumerate(prs.slide_masters):
        _collect_alt_texts(master.shapes, f"Master {i + 1}", alt_texts)
        for j, layout in enumerate(master.slide_layouts):
            _collect_alt_texts(layout.shapes, f"Layout {layout.name or j + 1}", alt_texts)

    return alt_texts


def _collect_alt_texts(shapes, label: str, result: list[tuple[str, str]]):
    """Recursively collect alt texts (descr attributes) from shapes."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                _collect_alt_texts(shape.shapes, label, result)
            except Exception:
                pass
            continue

        for elem in shape._element.iter():
            if etree.QName(elem).localname == "cNvPr":
                descr = (elem.get("descr") or "").strip()
                if descr:
                    result.append((descr, label))
                break  # Only one cNvPr per shape


def _replace_alt_texts_in_shapes_all(prs: Presentation, alt_mapping: dict[str, str]):
    """Replace alt texts across the entire presentation (slides, masters, layouts)."""
    for slide in prs.slides:
        _replace_alt_in_shapes(slide.shapes, alt_mapping)
    for master in prs.slide_masters:
        _replace_alt_in_shapes(master.shapes, alt_mapping)
        for layout in master.slide_layouts:
            _replace_alt_in_shapes(layout.shapes, alt_mapping)


def _replace_alt_in_shapes(shapes, alt_mapping: dict[str, str]):
    """Recursively replace alt texts in shapes."""
    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            try:
                _replace_alt_in_shapes(shape.shapes, alt_mapping)
            except Exception:
                pass
            continue

        for elem in shape._element.iter():
            if etree.QName(elem).localname == "cNvPr":
                descr = (elem.get("descr") or "").strip()
                if descr and descr in alt_mapping:
                    elem.set("descr", alt_mapping[descr])
                break


def _clean_metadata(prs: Presentation, entities_sorted: list[tuple[str, str]] | None = None):
    """Anonymize/clear metadata from presentation properties.

    The title is anonymized via entity replacement; other fields are cleared.
    """
    props = prs.core_properties

    # Anonymize title with entity replacement instead of just clearing
    if entities_sorted:
        try:
            title = props.title or ""
            if title:
                props.title = replace_entities_in_text(title, entities_sorted)
        except Exception:
            try:
                props.title = ""
            except Exception:
                pass
    else:
        try:
            props.title = ""
        except Exception:
            pass

    try:
        props.author = ""
    except Exception:
        pass
    try:
        props.last_modified_by = ""
    except Exception:
        pass
    try:
        props.comments = ""
    except Exception:
        pass
    try:
        props.keywords = ""
    except Exception:
        pass
    try:
        props.subject = ""
    except Exception:
        pass


def save_to_bytes(prs: Presentation) -> bytes:
    """Save a Presentation to bytes."""
    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()
